import ast
import base64
import ipaddress
import logging
import operator
import os
import re
import socket
import textwrap
from datetime import UTC, datetime
from urllib.parse import urlparse
from uuid import UUID
from zoneinfo import ZoneInfo

import httpx
from ddgs import DDGS
from langchain_core.runnables import RunnableConfig
from langchain_core.tools import tool
from sqlalchemy import select

from app.employees.models import Employee
from app.memory.service import recall, remember
from app.organizations.models import Organization

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# SSRF protection — hostname / IP blocklists
# ---------------------------------------------------------------------------

_BLOCKED_SCHEMES = {"file", "ftp", "gopher", "dict", "ldap", "tftp"}

_BLOCKED_HOSTNAMES = {
    "localhost",
    "127.0.0.1",
    "0.0.0.0",
    "::1",
    "[::1]",
    "0",
    "metadata.google.internal",  # cloud metadata endpoints
    "169.254.169.254",
}

_BLOCKED_NETWORKS = [
    ipaddress.ip_network("127.0.0.0/8"),      # loopback
    ipaddress.ip_network("10.0.0.0/8"),       # private
    ipaddress.ip_network("172.16.0.0/12"),    # private
    ipaddress.ip_network("192.168.0.0/16"),   # private
    ipaddress.ip_network("169.254.0.0/16"),   # link-local
    ipaddress.ip_network("0.0.0.0/8"),        # current network
    ipaddress.ip_network("::1/128"),          # IPv6 loopback
    ipaddress.ip_network("fc00::/7"),         # IPv6 unique local
    ipaddress.ip_network("fe80::/10"),        # IPv6 link-local
]


def _is_private_host(hostname: str) -> bool:
    """Return True if *hostname* resolves to or is a private / loopback address."""
    # Strip brackets from IPv6 addresses
    clean = hostname.strip("[]")
    # Fast path — literal IP
    try:
        addr = ipaddress.ip_address(clean)
        return any(addr in net for net in _BLOCKED_NETWORKS)
    except ValueError:
        pass  # not a literal IP, try DNS
    # DNS resolution
    try:
        resolved = socket.getaddrinfo(clean, None)
    except OSError:
        # Can't resolve — err on the side of safety
        return True
    for family, _, _, _, sockaddr in resolved:
        ip_str = sockaddr[0]
        try:
            addr = ipaddress.ip_address(ip_str)
        except ValueError:
            continue
        if any(addr in net for net in _BLOCKED_NETWORKS):
            return True
    return False


def _validate_url(url: str) -> str:
    """Validate URL is safe to fetch. Returns an error message or empty string."""
    parsed = urlparse(url)

    # Block dangerous schemes
    if parsed.scheme in _BLOCKED_SCHEMES:
        return f"Blocked scheme: {parsed.scheme}"

    if parsed.scheme not in ("http", "https"):
        return f"Unsupported scheme: {parsed.scheme or '(empty)'}"

    hostname = parsed.hostname
    if not hostname:
        return "No hostname in URL"

    # Block well-known blocked hostnames
    if hostname.lower() in _BLOCKED_HOSTNAMES or hostname in _BLOCKED_HOSTNAMES:
        return f"Blocked hostname: {hostname}"

    # Block cloud metadata IPs
    if hostname == "169.254.169.254":
        return f"Blocked hostname: {hostname}"

    # DNS / IP check
    if _is_private_host(hostname):
        return f"Blocked internal address: {hostname}"

    return ""  # safe


@tool
async def search_web(query: str) -> str:
    """Search the web for current information, news, or facts. Useful for time-sensitive queries."""
    try:
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=5))
        if not results:
            return "No results found."
        return "\n\n".join(
            f"**{r['title']}**\n{r['body']}\nSource: {r['href']}" for r in results
        )
    except Exception as e:
        return f"Error performing web search: {e}"


@tool
def get_datetime(timezone: str = "UTC") -> str:
    """Get the current date and time. Useful for scheduling and checking dates.

    Accepts an IANA timezone name (e.g. ``"America/New_York"``, ``"Asia/Tokyo"``).
    Defaults to UTC when the timezone is unrecognised.
    """
    tz_info: UTC | ZoneInfo
    try:
        tz_info = ZoneInfo(timezone)
    except Exception:
        tz_info = UTC
        timezone = "UTC"
    now = datetime.now(tz_info)
    return f"Current date and time: {now.strftime('%Y-%m-%d %H:%M:%S %Z')} ({timezone})"


@tool
def calculate(expression: str) -> str:
    """Safely evaluate a mathematical expression. E.g. '2 + 2 * 10 / 5'."""
    _ops = {
        ast.Add: operator.add,
        ast.Sub: operator.sub,
        ast.Mult: operator.mul,
        ast.Div: operator.truediv,
        ast.Pow: operator.pow,
    }

    def _eval(node):
        if isinstance(node, ast.Constant):
            return node.value
        if isinstance(node, ast.BinOp):
            return _ops[type(node.op)](_eval(node.left), _eval(node.right))
        if isinstance(node, ast.UnaryOp) and isinstance(node.op, ast.USub):
            return -_eval(node.operand)
        raise ValueError(f"Unsupported node type: {type(node).__name__}")

    try:
        # Parse expression safely
        tree = ast.parse(expression, mode="eval")
        result = _eval(tree.body)
        return str(result)
    except Exception as e:
        return f"Error evaluating expression: {e}"


@tool
async def fetch_url(url: str) -> str:
    """Fetch and return the text content of a public webpage. Blocks internal URLs."""
    # SSRF check before any network I/O
    err = _validate_url(url)
    if err:
        return f"Cannot fetch URL: {err}"

    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=10) as client:
            # Re-validate the final URL after redirects to catch open-redirect
            # to internal addresses. httpx doesn't expose the redirect chain
            # natively, so we issue a HEAD first to get the final URL.
            head = await client.head(url)
            final_url = str(head.url)
            final_err = _validate_url(final_url)
            if final_err:
                return f"Cannot fetch URL (redirect): {final_err}"

            resp = await client.get(url)
            resp.raise_for_status()

        # Simple HTML tag removal
        text = re.sub(r"<[^>]+>", " ", resp.text)
        text = re.sub(r"\s+", " ", text).strip()
        # Return first 3000 chars to avoid prompt bloat
        return text[:3000]
    except httpx.HTTPError as e:
        return f"Error fetching URL: {e}"
    except Exception as e:
        return f"Error fetching URL: {e}"


@tool
async def search_memory(query: str, config: RunnableConfig = None) -> str:
    """Search team memory for past decisions, facts, and knowledge.
    Searches both your personal memory and shared org knowledge."""
    emp_id = None
    db = None
    if config and "configurable" in config:
        emp_id = config["configurable"].get("employee_id")
        db = config["configurable"].get("db")

    if not emp_id or not db:
        return "Memory search unavailable (no employee context)."

    try:
        emp = await db.scalar(
            select(Employee).where(Employee.id == UUID(emp_id))
        )
        if not emp:
            return "Employee not found."

        org = await db.scalar(
            select(Organization).where(Organization.id == emp.org_id)
        )

        datasets: list[str] = []
        if emp.cognee_dataset_name:
            datasets.append(emp.cognee_dataset_name)
        if org and org.cognee_dataset_name:
            datasets.append(org.cognee_dataset_name)

        user_id = emp.cognee_user_id or (
            org.cognee_system_user_id if org else None
        )
        if not user_id or not datasets:
            return "Memory not yet provisioned for this employee."

        results = await recall(query, user_id, datasets=datasets)
        if not results:
            return f"No relevant memory found for '{query}'."

        lines = []
        for r in results:
            ds = r.get("dataset_name", "")
            text = r.get("text", "")
            lines.append(f"- [{ds}] {text}")
        return (
            f"Found {len(results)} memory result(s):\n\n"
            + "\n".join(lines)
        )
    except Exception as e:
        logger.exception("search_memory tool failed")
        return f"Error searching memory: {e}"


@tool
async def ingest_memory(content: str, config: RunnableConfig = None) -> str:
    """Store an important fact or decision in your personal memory
    for future reference. Use this to remember things you've learned."""
    emp_id = None
    db = None
    if config and "configurable" in config:
        emp_id = config["configurable"].get("employee_id")
        db = config["configurable"].get("db")

    if not emp_id or not db:
        return "Cannot store memory (no employee context)."

    try:
        emp = await db.scalar(
            select(Employee).where(Employee.id == UUID(emp_id))
        )
        if (
            not emp
            or not emp.cognee_user_id
            or not emp.cognee_dataset_name
        ):
            return "Memory not yet provisioned for this employee."

        await remember(
            content,
            emp.cognee_dataset_name,
            emp.cognee_user_id,
            dataset_id=emp.cognee_dataset_id,
            background=False,
        )
        return (
            f"Successfully remembered: {content[:200]}"
            + ("..." if len(content) > 200 else "")
        )
    except Exception as e:
        logger.exception("ingest_memory tool failed")
        return f"Error storing memory: {e}"


# ---------------------------------------------------------------------------
# File generation — marker-based: tool returns a marker string and the
# formatter node picks it up and populates state["files"].
# ---------------------------------------------------------------------------

_FILE_MARKER_PREFIX = "__OPENHUMAN_FILE__"


def _file_marker(filename: str, content_type: str, data_b64: str, title: str = "") -> str:
    """Build a marker string the formatter will intercept as a file upload."""
    import json
    payload = json.dumps({
        "filename": filename,
        "content_type": content_type,
        "data": data_b64,
        "title": title or filename,
    })
    return f"{_FILE_MARKER_PREFIX}{payload}"


def _find_dejavu_font(style: str = "") -> str:
    filenames = {"": "DejaVuSans.ttf", "B": "DejaVuSans-Bold.ttf"}
    filename = filenames.get(style, f"DejaVuSans{style}.ttf")
    search_paths = [
        "/usr/share/fonts/truetype/dejavu",
        "/usr/share/fonts/TTF",
        "/usr/share/fonts/dejavu",
        "/usr/local/share/fonts",
    ]
    for base in search_paths:
        path = os.path.join(base, filename)
        if os.path.exists(path):
            return path
    raise FileNotFoundError(f"DejaVu font file not found: {filename}")


def _generate_pdf(content: str) -> bytes:
    from fpdf import FPDF

    def _write_wrapped_line(pdf: FPDF, text: str, line_height: int = 6, indent: int = 0) -> None:
        clean = text.replace("\t", "    ").strip()
        if not clean:
            pdf.ln(4)
            return

        available_width = max(pdf.w - pdf.l_margin - pdf.r_margin - indent, 20)
        chunks = textwrap.wrap(
            clean,
            width=90,
            break_long_words=True,
            break_on_hyphens=False,
            replace_whitespace=False,
            drop_whitespace=False,
        ) or [clean]

        for chunk in chunks:
            if indent:
                pdf.set_x(pdf.l_margin + indent)
            pdf.multi_cell(available_width, line_height, chunk)

    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.add_font("DejaVu", "", _find_dejavu_font(""), uni=True)
    pdf.add_font("DejaVu", "B", _find_dejavu_font("B"), uni=True)
    pdf.set_font("DejaVu", size=11)
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            pdf.ln(4)
            continue
        if line.startswith("---"):
            continue
        if line.startswith("##"):
            pdf.set_font("DejaVu", "B", 14)
            _write_wrapped_line(pdf, line.lstrip("#").strip(), line_height=8)
            pdf.set_font("DejaVu", size=11)
        elif line.startswith("#"):
            pdf.set_font("DejaVu", "B", 16)
            _write_wrapped_line(pdf, line.lstrip("#").strip(), line_height=10)
            pdf.set_font("DejaVu", size=11)
        elif any(line.startswith(p) for p in ("- ", "* ", "•")):
            _write_wrapped_line(pdf, f"• {line.lstrip('- *•').strip()}", indent=5)
        elif line[0].isdigit() and ". " in line[:4]:
            _write_wrapped_line(pdf, line, indent=5)
        else:
            _write_wrapped_line(pdf, line)
    return pdf.output()


def _generate_pptx(content: str) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    bg_color = RGBColor(0x1A, 0x1A, 0x2E)
    text_color = RGBColor(0xF0, 0xF0, 0xFF)
    accent = RGBColor(0x6C, 0x63, 0xFF)

    lines = content.strip().split("\n")
    slide_content: list[str] = []
    title = ""

    def _flush_slide():
        nonlocal title
        if not slide_content and not title:
            return
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        if title:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.color.rgb = accent
            p.font.bold = True
        y_start = Inches(1.6) if title else Inches(0.5)
        txBox = slide.shapes.add_textbox(Inches(0.8), y_start, Inches(11.7), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(slide_content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = text_color
            p.space_after = Pt(6)
        slide_content.clear()
        title = ""

    for line in lines:
        line = line.strip()
        if not line or line.startswith("---"):
            _flush_slide()
            continue
        if line.startswith("##") or line.startswith("###"):
            _flush_slide()
            title = line.lstrip("#").strip()
        elif line.startswith("#"):
            _flush_slide()
            title = line.lstrip("#").strip()
        elif any(line.startswith(p) for p in ("- ", "* ", "•")):
            slide_content.append(line.lstrip("- *•").strip())
        elif line[0].isdigit() and ". " in line[:4]:
            slide_content.append(line)
        else:
            slide_content.append(line)
    _flush_slide()
    from io import BytesIO
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


@tool
def create_document(content: str, filename: str = "document.txt") -> str:
    """Create a downloadable document from content.

    Generates a proper file (PDF, PPTX, or plain text) from the given
    content. Use this after fetching or generating content that the user
    may want to save — e.g. reports, pitch deck outlines, summaries.

    The file will be uploaded to the conversation as an attachment.
    For pitch decks / presentations use a ``.pptx`` extension.
    For formal reports use ``.pdf``.
    For quick notes use ``.txt`` or ``.md``.

    Args:
        content: The full text content.
        filename: Desired filename with extension
                  (e.g. "report.pdf", "deck.pptx", "notes.md").
    """
    import os

    ext = os.path.splitext(filename)[1].lower()

    try:
        if ext == ".pdf":
            raw = _generate_pdf(content)
            data_b64 = base64.b64encode(raw).decode("ascii")
            ct = "application/pdf"
        elif ext == ".pptx":
            raw = _generate_pptx(content)
            data_b64 = base64.b64encode(raw).decode("ascii")
            ct = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        else:
            data_b64 = base64.b64encode(content.encode("utf-8")).decode("ascii")
            from mimetypes import guess_type
            ct, _ = guess_type(filename) if "." in filename else ("text/plain", None)
            ct = ct or "text/plain"
    except Exception:
        logger.exception("Document generation failed for %s", filename)
        fallback_name = re.sub(r"\.[^.]+$", "", filename) + ".txt"
        data_b64 = base64.b64encode(content.encode("utf-8")).decode("ascii")
        ct = "text/plain"
        return _file_marker(
            fallback_name,
            ct,
            data_b64,
            title=f"{fallback_name} (fallback after document generation error)",
        )

    return _file_marker(filename, ct, data_b64, title=filename)


from app.agent.tools.cancel_background_task import cancel_background_task  # noqa: E402
from app.agent.tools.check_background_task import check_background_task  # noqa: E402
from app.agent.tools.cronjob_tools import cronjob  # noqa: E402
from app.agent.tools.delegate_task import delegate_task  # noqa: E402
from app.agent.tools.escalation import (
    escalate_to_human,  # noqa: E402
    escalate_to_human_interactive,  # noqa: E402
)
from app.agent.tools.vision_tools import vision_analyze  # noqa: E402
from app.agent.tools.file_tools import (  # noqa: E402
    file_read,
    file_write,
    file_search,
    file_patch,
    file_delete,
)
from app.agent.tools.todo_tool import todo  # noqa: E402
from app.agent.tools.code_execution_tool import execute_code  # noqa: E402
from app.agent.tools.image_generation_tool import generate_image  # noqa: E402
from app.agent.tools.transcription_tools import transcribe_audio  # noqa: E402

# List of all built-in tools to export
BUILT_IN_TOOLS = [
    search_web,
    get_datetime,
    calculate,
    fetch_url,
    search_memory,
    ingest_memory,
    create_document,
    check_background_task,
    cancel_background_task,
    escalate_to_human,
    escalate_to_human_interactive,
    delegate_task,
    cronjob,
    vision_analyze,
    file_read,
    file_write,
    file_search,
    file_patch,
    file_delete,
    todo,
    execute_code,
    generate_image,
    transcribe_audio,
]
